from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x1A, 0x56, 0xDB)
LIGHT_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x1E, 0x1E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x16, 0x7F, 0x3D)
LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)
PEACH = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_PEACH = RGBColor(0xFD, 0xED, 0xED)
MAUVE = RGBColor(0x6C, 0x3C, 0xB0)
YELLOW = RGBColor(0xB8, 0x86, 0x0B)
FONT_NAME = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_divider(slide, left, top, width, color=ACCENT):
    """Add a horizontal accent line under a slide title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle card as a visual container."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_icon_circle(slide, left, top, size, fill_color, label, font_color=WHITE):
    """Add a colored circle with a short label as a visual icon."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(int(size / Inches(1) * 11))
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
# Accent bar at top
add_rounded_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Employee Records System", font_size=44, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_divider(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT)
add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8),
             "Python CLI  |  SQLite Database  |  Matplotlib Charts",
             font_size=22, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
             "Belvah Shanyisa  •  April 2026",
             font_size=18, color=MAUVE, alignment=PP_ALIGN.CENTER)

# ── Slide 2: Problem Statement ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Problem Statement", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT)
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
    "• Track weekly ticket completions per employee",
    "• Manual tracking is error-prone, hard to visualize",
    "• Need quick views by department or whole team",
    "• Rank metrics highest to lowest",
    "• Persistent storage across sessions",
], font_size=20)

# ── Slide 3: Features ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Features", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT)
# Left card — CRUD
add_rounded_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), LIGHT_ACCENT, ACCENT)
add_icon_circle(slide, Inches(1.2), Inches(1.9), Inches(0.6), ACCENT, "C")
add_text_box(slide, Inches(2.0), Inches(1.95), Inches(4), Inches(0.5),
             "CRUD Operations", font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.5), [
    "• View all records ranked by tickets",
    "• Filter by department",
    "• Search by name (partial match)",
    "• Add with validation",
    "• Update by ID",
    "• Delete with confirmation",
], font_size=17, color=DARK_GRAY)
# Right card — Visualization
add_rounded_box(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8), LIGHT_ACCENT, ACCENT)
add_icon_circle(slide, Inches(7.2), Inches(1.9), Inches(0.6), MAUVE, "V")
add_text_box(slide, Inches(8.0), Inches(1.95), Inches(4), Inches(0.5),
             "Visualization & Validation", font_size=20, color=MAUVE, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(2.7), Inches(5), Inches(3.5), [
    "• Terminal bar chart (Unicode)",
    "• Matplotlib graphical chart",
    "• Unique color shade per employee",
    "• Duplicate prevention",
    "• Input validation (alphanumeric, non-negative)",
], font_size=17, color=DARK_GRAY)

# ── Slide 4: Bar Chart Visualization ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Bar Chart Visualization", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(3.5), ACCENT)
slide.shapes.add_picture(
    "/Users/shanyisa/Desktop/Screenshot 2026-04-20 at 00.23.59.png",
    Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.5))

# ── Slide 5: Tech Stack ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Tech Stack", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2.2), ACCENT)
# Row of icon cards
card_data = [
    (Inches(0.8),  ACCENT, "Py", "Python 3", "Core language"),
    (Inches(3.9),  GREEN,  "SQ", "SQLite", "Persistent database"),
    (Inches(7.0),  MAUVE,  "Mp", "matplotlib", "Graphical charts"),
    (Inches(10.1), YELLOW, "Cs", "colorsys", "HSL color palettes"),
]
for left, icon_color, icon_label, title, desc in card_data:
    add_rounded_box(slide, left, Inches(1.8), Inches(2.6), Inches(3.2), LIGHT_ACCENT, ACCENT)
    add_icon_circle(slide, left + Inches(0.85), Inches(2.1), Inches(0.9), icon_color, icon_label)
    add_text_box(slide, left + Inches(0.1), Inches(3.3), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.9), Inches(2.4), Inches(0.5),
                 desc, font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
# unittest card row
add_rounded_box(slide, Inches(3.9), Inches(5.3), Inches(5.5), Inches(1.2), LIGHT_ACCENT, ACCENT)
add_text_box(slide, Inches(4.1), Inches(5.5), Inches(5.1), Inches(0.4),
             "unittest — 33 automated tests", font_size=16, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.1), Inches(5.95), Inches(5.1), Inches(0.4),
           "Built-in Python testing framework", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 6: Code Structure (Modularized) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Code Structure — Modularized", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(4.5), ACCENT)
# Three module cards
# database.py
add_rounded_box(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.2), LIGHT_ACCENT, ACCENT)
add_icon_circle(slide, Inches(1.0), Inches(1.9), Inches(0.55), GREEN, "DB")
add_text_box(slide, Inches(1.7), Inches(1.95), Inches(2.5), Inches(0.5),
             "database.py", font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, Inches(0.9), Inches(2.7), Inches(3.3), Inches(4), [
    "• SQLite init + seeding",
    "• get_all / get_by_id",
    "• get_by_department",
    "• search_by_name",
    "• create / update / delete",
    "• has_duplicate",
], font_size=15, color=DARK_GRAY)
# charts.py
add_rounded_box(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(5.2), LIGHT_ACCENT, MAUVE)
add_icon_circle(slide, Inches(5.2), Inches(1.9), Inches(0.55), MAUVE, "Ch")
add_text_box(slide, Inches(5.9), Inches(1.95), Inches(2.5), Inches(0.5),
             "charts.py", font_size=18, color=MAUVE, bold=True)
add_bullet_list(slide, Inches(5.1), Inches(2.7), Inches(3.3), Inches(4), [
    "• Terminal Unicode bar chart",
    "• Matplotlib graphical chart",
    "• generate_palette()",
    "• get_employee_color()",
    "• Dept-based color schemes",
], font_size=15, color=DARK_GRAY)
# employees_records_system.py
add_rounded_box(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(5.2), LIGHT_ACCENT, ACCENT)
add_icon_circle(slide, Inches(9.4), Inches(1.9), Inches(0.55), ACCENT, "UI")
add_text_box(slide, Inches(10.1), Inches(1.95), Inches(2.5), Inches(0.5),
             "main app", font_size=18, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(9.3), Inches(2.7), Inches(3.3), Inches(4), [
    "• CLI menu loop",
    "• User input handling",
    "• display_records()",
    "• view_all / view_by_dept",
    "• create / update / delete",
    "• search_employee",
], font_size=15, color=DARK_GRAY)

# ── Slide 7: Key Concepts Used ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Key Python Concepts", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(3.5), ACCENT)
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), [
    "• Lists, dictionaries, SQL queries",
    "• Functions with docstrings",
    "• Modular imports across files",
    "• While loops and conditionals",
    "• f-string formatting",
], font_size=20)
add_bullet_list(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(5), [
    "• try/except error handling",
    "• List comprehensions",
    "• Lambda sorting",
    "• SQLite parameterized queries",
    "• Context management (try/finally)",
], font_size=20)

# ── Slide 8: Validation Rules ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Validation Rules", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2.8), ACCENT)
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
    "• Alphanumeric usernames only",
    "• Non-empty username required",
    "• Department must match predefined list",
    "• No duplicate name + department",
    "• Non-negative integer tickets",
    "• try/except prevents crashes",
], font_size=20)

# ── Slide 9: Demo – Main Menu ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Demo — Main Menu", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT)
slide.shapes.add_picture(
    "/Users/shanyisa/Desktop/Screenshot 2026-04-20 at 00.23.23.png",
    Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.5))

# ── Slide 10: Demo – Search by Name ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Demo — Search by Name", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(3), ACCENT)
slide.shapes.add_picture(
    "/Users/shanyisa/Desktop/Screenshot 2026-04-20 at 00.25.17.png",
    Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.5))

# ── Slide 11: Demo – Bar Chart ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Demo — Bar Chart Visualization", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(4.5), ACCENT)
slide.shapes.add_picture(
    "/Users/shanyisa/Desktop/Screenshot 2026-04-20 at 00.25.41.png",
    Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.5))

# ── Slide 12: Challenges & Lessons Learned ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Challenges & Lessons Learned", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(4.5), ACCENT)
# Challenges card
add_rounded_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5), LIGHT_PEACH, PEACH)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(4.5), Inches(0.5),
             "Challenges", font_size=22, color=PEACH, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.5), Inches(4.9), Inches(3.2), [
    "• Unique color shades per employee",
    "• Unicode bar chart alignment",
    "• Edge cases in user input",
    "• Organizing code across modules",
], font_size=17, color=DARK_GRAY)
# Lessons card
add_rounded_box(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.5), LIGHT_GREEN, GREEN)
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(4.5), Inches(0.5),
             "Lessons Learned", font_size=22, color=GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.2), [
    "• Modular code = easier maintenance",
    "• Input validation prevents crashes",
    "• SQLite > in-memory for persistence",
    "• Unit testing catches bugs early",
], font_size=17, color=DARK_GRAY)

# ── Slide 13: Future Improvements ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Future Improvements", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(3.2), ACCENT)
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
    "• Export to CSV, PDF, Excel",
    "• Configurable departments",
    "• Web interface (Flask / Django)",
], font_size=20)

# ── Slide 14: Unit Testing ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Unit Testing", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT)
# Summary badge
add_rounded_box(slide, Inches(4.0), Inches(1.6), Inches(5.3), Inches(0.7), GREEN, GREEN)
add_text_box(slide, Inches(4.2), Inches(1.65), Inches(5), Inches(0.6),
             "33 tests  •  All passing", font_size=20, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(4.5), [
    "• Database init + seeding",
    "• Color palettes — hex, unique",
    "• Department filtering",
    "• Sorting by tickets",
    "• Duplicate detection",
], font_size=18)
add_bullet_list(slide, Inches(6.5), Inches(2.7), Inches(6), Inches(4.5), [
    "• Input validation",
    "• Search — exact, partial, case",
    "• Create + delete via DB",
    "• Update + nonexistent handling",
    "• Isolated test database",
], font_size=18)

# ── Slide 15: Conclusion ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Conclusion", font_size=36, color=ACCENT, bold=True)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2.2), ACCENT)
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5), [
    "• Full CRUD CLI app with SQLite persistence",
    "• Dual visualizations — terminal + matplotlib",
    "• Modularized: database, charts, main app",
    "• Docstrings on every function",
    "• 33 passing unit tests",
], font_size=20)

# ── Slide 16: Q&A ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE_BG)
add_rounded_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Questions?", font_size=44, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_divider(slide, Inches(5), Inches(3.7), Inches(3.3), ACCENT)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "GitHub: github.com/Belvah", font_size=20, color=DARK_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6),
             "Belvah Shanyisa  •  misshuey3@gmail.com", font_size=18,
             color=MAUVE, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/shanyisa/Desktop/python_foundations/python_intro/employee_weekly/Employee_Records_System_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
